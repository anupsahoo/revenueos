"""Builds docs/slides/RevenueOS.pptx.

Ten slides, same purple palette as the operator screen. Regenerate with:

    python3 -m venv .venv && .venv/bin/pip install python-pptx
    .venv/bin/python docs/slides/build_deck.py

Kept in the repo so the deck is reproducible rather than hand-edited.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# palette — matches app/globals.css
INK = RGBColor(0x1C, 0x16, 0x2E)
MUTED = RGBColor(0x6B, 0x64, 0x82)
ACCENT = RGBColor(0x6D, 0x28, 0xD9)
ACCENT_SOFT = RGBColor(0xED, 0xE7, 0xFB)
ACCENT_DEEP = RGBColor(0x3B, 0x0F, 0x76)
PAPER = RGBColor(0xFA, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x15, 0x80, 0x3D)
AMBER = RGBColor(0xB4, 0x53, 0x09)
RED = RGBColor(0xB9, 0x1C, 0x1C)

W, H = Inches(13.333), Inches(7.5)


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (string, size_pt, color, bold, space_after_pt)"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (t, size, color, bold, after) in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(after)
        r = para.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Helvetica Neue"
    return box


def box(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.25)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, KeyError):
            pass
    shape.text_frame.word_wrap = True
    return shape


def label(shape, s, size=14, color=INK, bold=False):
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Helvetica Neue"


def heading(slide, title, kicker=None):
    box(slide, Inches(0.7), Inches(0.62), Inches(0.09), Inches(0.52), ACCENT, radius=False)
    text(slide, Inches(0.95), Inches(0.5), Inches(11.5), Inches(0.8),
         [(title, 30, INK, True, 0)])
    if kicker:
        text(slide, Inches(0.95), Inches(1.18), Inches(11.5), Inches(0.4),
             [(kicker, 14, MUTED, False, 0)])


def arrow(slide, x, y, w, h=Inches(0.2), color=ACCENT):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def footer(slide, n):
    text(slide, Inches(11.9), Inches(6.85), Inches(0.9), Inches(0.3),
         [(str(n), 11, MUTED, False, 0)], align=PP_ALIGN.RIGHT)


def build(out: Path):
    prs = deck()

    # 1 — title
    s = blank(prs, ACCENT_DEEP)
    box(s, Inches(0), Inches(0), Inches(0.35), H, ACCENT, radius=False)
    text(s, Inches(1.2), Inches(2.3), Inches(11), Inches(1.4),
         [("RevenueOS", 60, WHITE, True, 8),
          ("One seam, closed end to end.", 26, RGBColor(0xC9, 0xB8, 0xF5), False, 0)])
    text(s, Inches(1.2), Inches(4.3), Inches(11), Inches(1.6),
         [("Sales wins a deal. PreSales has two business days to start the POC plan.",
           17, RGBColor(0xDD, 0xD3, 0xF8), False, 6),
          ("In the US it takes 6.8. This is the loop that closes the gap.",
           17, RGBColor(0xDD, 0xD3, 0xF8), False, 0)])

    # 2 — the seam, four numbers
    s = blank(prs)
    heading(s, "The seam", "Marketing → Sales → PreSales → Delivery → Support. It stalls at one join.")
    chain = ["Marketing", "Sales", "PreSales", "Delivery", "Support"]
    x, y, bw, bh = Inches(0.95), Inches(1.95), Inches(2.05), Inches(0.75)
    for i, name in enumerate(chain):
        stalled = name == "PreSales"
        b = box(s, x + Inches(2.4) * i, y, bw, bh,
                ACCENT if stalled else WHITE, None if stalled else RGBColor(0xDD, 0xD8, 0xEA))
        label(b, name, 14, WHITE if stalled else INK, stalled)
        if i < 4:
            arrow(s, x + Inches(2.4) * i + bw + Inches(0.06), y + Inches(0.28), Inches(0.24))
    box(s, x + Inches(4.8), y + bh + Inches(0.12), bw, Inches(0.06), RED, radius=False)
    text(s, x + Inches(4.55), y + bh + Inches(0.24), Inches(2.6), Inches(0.3),
         [("the stall", 12, RED, True, 0)], align=PP_ALIGN.CENTER)

    nums = [("2 days", "the SLA", ACCENT), ("6.8 days", "US actual", RED),
            ("21%", "US reuse", RED), ("58% / 55%", "UK / India reuse", GREEN)]
    for i, (big, cap, col) in enumerate(nums):
        b = box(s, Inches(0.95) + Inches(3.0) * i, Inches(3.6), Inches(2.7), Inches(1.5),
                WHITE, RGBColor(0xDD, 0xD8, 0xEA))
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r = p1.add_run(); r.text = big
        r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = col
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = cap
        r2.font.size = Pt(13); r2.font.color.rgb = MUTED
    text(s, Inches(0.95), Inches(5.5), Inches(11.5), Inches(0.9),
         [("Not a headcount problem. The templates already exist — nothing fetches them.",
           19, INK, True, 0)])
    footer(s, 2)

    # 3 — the loop
    s = blank(prs)
    heading(s, "The loop", "Five steps. One screen. The human presses the last button.")
    steps = [("Brief\narrives", ACCENT_SOFT, INK), ("Retrieve\n3 matches", ACCENT_SOFT, INK),
             ("Draft\nPOC plan", ACCENT_SOFT, INK), ("Human\naccepts", ACCENT, WHITE),
             ("Library\nre-ranks", ACCENT_SOFT, INK)]
    x, y, bw, bh = Inches(0.95), Inches(2.35), Inches(2.05), Inches(1.5)
    for i, (name, fill, col) in enumerate(steps):
        b = box(s, x + Inches(2.4) * i, y, bw, bh, fill,
                ACCENT if col == INK else None)
        label(b, name, 15, col, True)
        if i < 4:
            arrow(s, x + Inches(2.4) * i + bw + Inches(0.06), y + Inches(0.65), Inches(0.24))
    curve = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, x + Inches(0.9), y + bh + Inches(0.45),
                               Inches(9.0), Inches(0.22))
    curve.fill.solid(); curve.fill.fore_color.rgb = RGBColor(0xB9, 0xA6, 0xEE)
    curve.line.fill.background(); curve.shadow.inherit = False
    text(s, x + Inches(2.6), y + bh + Inches(0.72), Inches(6.0), Inches(0.4),
         [("every decision changes the next retrieval", 13, MUTED, False, 0)],
         align=PP_ALIGN.CENTER)
    text(s, Inches(0.95), Inches(5.85), Inches(11.5), Inches(0.7),
         [("Automated up to the commitment. The commitment stays with a named human.",
           19, INK, True, 0)])
    footer(s, 3)

    # 4 — event log
    s = blank(prs)
    heading(s, "One source of truth", "Nothing on the screen is stored. Every number is derived from an append-only log.")
    b = box(s, Inches(0.95), Inches(2.0), Inches(3.6), Inches(3.4), ACCENT, None)
    label(b, "Event log\n\nbrief.arrived\ndraft.generated\ndraft.accepted\ndraft.rejected\nsla.breached\ntrigger.fired",
          14, WHITE, False)
    arrow(s, Inches(4.75), Inches(3.55), Inches(0.5))
    derived = ["age → business days since arrival",
               "status → green / amber / red",
               "reuse rate → accepted ÷ decided",
               "boosts → ±8 per decision",
               "triggers → who was escalated"]
    for i, d in enumerate(derived):
        db = box(s, Inches(5.5), Inches(2.0) + Inches(0.68) * i, Inches(6.9), Inches(0.55),
                 WHITE, RGBColor(0xDD, 0xD8, 0xEA))
        tf = db.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        p.text = "  " + d
        p.runs[0].font.size = Pt(14); p.runs[0].font.color.rgb = INK
    text(s, Inches(0.95), Inches(5.75), Inches(11.5), Inches(0.7),
         [("Every number on screen has an ⓘ that shows the query behind it. "
           "If it cannot, it does not go on the screen.", 17, ACCENT, True, 0)])
    footer(s, 4)

    # 5 — retrieval
    s = blank(prs)
    heading(s, "Retrieval you can argue with", "A score a Solution Architect can push back on beats a similarity number they cannot.")
    rows = [("Same segment", "30"), ("Same regulator / analogue", "16 / 8"),
            ("Problem overlap", "≤ 26"), ("Systems overlap", "≤ 18"),
            ("Recently used", "5"), ("Proven reuse", "4"),
            ("Learned from decisions", "± 8")]
    for i, (name, pts) in enumerate(rows):
        yy = Inches(2.05) + Inches(0.62) * i
        rb = box(s, Inches(0.95), yy, Inches(5.6), Inches(0.5), WHITE, RGBColor(0xDD, 0xD8, 0xEA))
        tf = rb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = "  " + name
        p.runs[0].font.size = Pt(14); p.runs[0].font.color.rgb = INK
        pb = box(s, Inches(6.65), yy, Inches(1.0), Inches(0.5), ACCENT_SOFT, ACCENT)
        label(pb, pts, 14, ACCENT_DEEP, True)
    tb = box(s, Inches(8.4), Inches(2.05), Inches(4.0), Inches(2.4), ACCENT_SOFT, ACCENT)
    label(tb, "Threshold 40\n\nNothing below it\nis shown.\nNo match = the plan\nsays so.", 16, ACCENT_DEEP, False)
    text(s, Inches(8.4), Inches(4.75), Inches(4.0), Inches(1.6),
         [("Each match ships the reasons it fits, in words.", 15, INK, True, 6),
          ("At ~200 templates I put embeddings in front for recall and keep this for the explanation.",
           13, MUTED, False, 0)])
    footer(s, 5)

    # 6 — the agent
    s = blank(prs)
    heading(s, "The agent", "Three nodes, one job. It drafts; it does not decide.")
    nodes = [("retrieve", "score the library,\nkeep the reasons"),
             ("draft", "write the POC plan\nfrom the matches"),
             ("assemble", "plan + handoff +\nwhat it used")]
    for i, (name, desc) in enumerate(nodes):
        nb = box(s, Inches(0.95) + Inches(4.1) * i, Inches(2.2), Inches(3.5), Inches(1.9),
                 WHITE, ACCENT)
        tf = nb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = name
        r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = ACCENT
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = "\n" + desc
        r2.font.size = Pt(13); r2.font.color.rgb = MUTED
        if i < 2:
            arrow(s, Inches(0.95) + Inches(4.1) * i + Inches(3.56), Inches(3.05), Inches(0.44))
    gb = box(s, Inches(0.95), Inches(4.5), Inches(11.45), Inches(1.0), ACCENT_SOFT, ACCENT)
    label(gb, "If the model is slow or unavailable, it falls back to a deterministic plan — "
              "and the screen says “sample draft” instead of “drafted by AI model”.", 15, ACCENT_DEEP)
    text(s, Inches(0.95), Inches(5.8), Inches(11.5), Inches(0.6),
         [("The demo never dies, and it never pretends.", 18, INK, True, 0)])
    footer(s, 6)

    # 7 — ask the seam
    s = blank(prs)
    heading(s, "Ask the seam", "A chat box with no memory and six read-only tools.")
    tools = ["get_events", "get_health", "get_brief", "get_template", "explain_match", "read_doc"]
    for i, t in enumerate(tools):
        tb = box(s, Inches(0.95) + Inches(3.9) * (i % 3), Inches(2.1) + Inches(0.85) * (i // 3),
                 Inches(3.5), Inches(0.65), WHITE, ACCENT)
        label(tb, t, 15, ACCENT_DEEP, True)
    qb = box(s, Inches(0.95), Inches(4.1), Inches(11.45), Inches(1.05), ACCENT, None)
    label(qb, "“Which briefs are past SLA right now and who was triggered?”  →  "
              "answers from the log, and names the sources it read.", 15, WHITE)
    xb = box(s, Inches(0.95), Inches(5.35), Inches(11.45), Inches(0.85), WHITE, RGBColor(0xDD, 0xD8, 0xEA))
    label(xb, "“What is the weather in Mumbai?”  →  refuses in one sentence. "
              "No tool holds it, so there is no answer.", 15, MUTED)
    footer(s, 7)

    # 8 — what it learns
    s = blank(prs)
    heading(s, "It learns from the decision, not from a rating", "No feedback form. The button is the signal.")
    a = box(s, Inches(0.95), Inches(2.2), Inches(5.4), Inches(1.5), WHITE, GREEN)
    label(a, "Accept  →  + 8 to every template that draft used", 18, GREEN, True)
    r = box(s, Inches(7.0), Inches(2.2), Inches(5.4), Inches(1.5), WHITE, RED)
    label(r, "Reject  →  − 8 to every template that draft used", 18, RED, True)
    arrow(s, Inches(6.55), Inches(4.0), Inches(0.28), Inches(0.2), MUTED)
    lb = box(s, Inches(0.95), Inches(4.15), Inches(11.45), Inches(1.15), ACCENT_SOFT, ACCENT)
    label(lb, "Boosts are derived from the event log — the server reads them itself "
              "and ignores anything the browser sends.", 16, ACCENT_DEEP)
    text(s, Inches(0.95), Inches(5.6), Inches(11.5), Inches(0.7),
         [("Refresh the page and the decision is still there. It was never browser state.",
           18, INK, True, 0)])
    footer(s, 8)

    # 9 — honesty
    s = blank(prs)
    heading(s, "What is real, and what breaks first", "Said before anyone has to ask.")
    left = box(s, Inches(0.95), Inches(2.0), Inches(5.6), Inches(3.9), WHITE, RGBColor(0xDD, 0xD8, 0xEA))
    tf = left.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25)
    for i, (t, sz, c, b) in enumerate([
        ("Real", 20, GREEN, True),
        ("Event log, derive functions, tests", 15, INK, False),
        ("Retrieval scoring and its reasons", 15, INK, False),
        ("The agent, the tools, the refusals", 15, INK, False),
        ("SLA breach detection and escalation", 15, INK, False),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run(); run.text = t
        run.font.size = Pt(sz); run.font.bold = b; run.font.color.rgb = c
    right = box(s, Inches(6.8), Inches(2.0), Inches(5.6), Inches(3.9), WHITE, RGBColor(0xDD, 0xD8, 0xEA))
    tf = right.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25)
    for i, (t, sz, c, b) in enumerate([
        ("Not real yet", 20, AMBER, True),
        ("Data is synthetic and labelled on screen", 15, INK, False),
        ("Store is in memory — a cold start loses it", 15, INK, False),
        ("No auth, one tenant, no workers", 15, INK, False),
        ("Breaks first at 10×: the store (#53)", 15, INK, False),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run(); run.text = t
        run.font.size = Pt(sz); run.font.bold = b; run.font.color.rgb = c
    text(s, Inches(0.95), Inches(6.15), Inches(11.5), Inches(0.6),
         [("48 open issues across six milestones. Nothing is closed unless code exists behind it.",
           16, MUTED, False, 0)])
    footer(s, 9)

    # 10 — the decision
    s = blank(prs, ACCENT_DEEP)
    box(s, Inches(0), Inches(0), Inches(0.35), H, ACCENT, radius=False)
    text(s, Inches(1.2), Inches(1.5), Inches(11), Inches(0.6),
         [("The one thing I would never hand to the agent", 30, WHITE, True, 0)])
    text(s, Inches(1.2), Inches(2.6), Inches(10.8), Inches(3.0),
         [("The accept.", 44, RGBColor(0xC9, 0xB8, 0xF5), True, 18),
          ("Pressing Accept is not a document being finished. It is a promise to a customer "
           "about what we will prove, in how many weeks, under their regulator.",
           19, RGBColor(0xDD, 0xD3, 0xF8), False, 12),
          ("The person who makes that promise has to be the person who will be in the room "
           "when it slips.", 19, RGBColor(0xDD, 0xD3, 0xF8), False, 0)])
    text(s, Inches(1.2), Inches(6.4), Inches(11), Inches(0.4),
         [("github.com/anupsahoo/revenueos   ·   revenueos-blond.vercel.app",
           14, RGBColor(0xB9, 0xA6, 0xEE), False, 0)])

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    build(Path(__file__).resolve().parent / "RevenueOS.pptx")
